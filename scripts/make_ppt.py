"""
6주차 발표용 PPT 생성 스크립트
스타일: D:/PARA/Assets/Templates/capstone_ppt_style.md 기준
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches
import copy
from lxml import etree
import os

# ── 색상 ──────────────────────────────────────────────
BG      = RGBColor(0xFA, 0xFA, 0xFA)  # 배경
NAVY    = RGBColor(0x0E, 0x28, 0x41)  # 딥 네이비
BLUE    = RGBColor(0x15, 0x60, 0x82)  # 메인 블루
SKYBLUE = RGBColor(0x0F, 0x9E, 0xD5)  # 밝은 블루
ORANGE  = RGBColor(0xE9, 0x71, 0x32)  # 포인트 주황
LGRAY   = RGBColor(0xE8, 0xE8, 0xE8)  # 밝은 회색
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY   = RGBColor(0x59, 0x59, 0x59)

# ── 폰트 ──────────────────────────────────────────────
F_TITLE  = 'A2Z 3 Light'
F_BOLD   = 'A2Z 6 SemiBold'
F_MED    = 'A2Z 5 Medium'
F_BODY   = 'Gowun Batang'

# ── 슬라이드 크기 ──────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)

# ── 여백 ──────────────────────────────────────────────
ML = Cm(1.2)   # left margin
MT = Cm(1.0)   # top margin
CW = W - ML * 2  # content width


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def set_bg(slide, color: RGBColor):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h,
                font_name=F_BODY, font_size=20, bold=False,
                color=NAVY, align=PP_ALIGN.LEFT,
                line_spacing=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    if line_spacing:
        p.line_spacing = line_spacing
    return txb


def add_multiline_textbox(slide, lines, x, y, w, h,
                          font_name=F_BODY, font_size=18,
                          bold=False, color=NAVY,
                          align=PP_ALIGN.LEFT,
                          highlight_pairs=None):
    """여러 줄 텍스트박스. highlight_pairs: [(line_idx, word, color), ...]"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name  = font_name
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return txb


def add_slide_number(slide, num, total):
    """우하단 슬라이드 번호"""
    add_textbox(slide, f'{num} / {total}',
                W - Cm(3.5), H - Cm(1.0), Cm(3.0), Cm(0.8),
                font_name=F_BODY, font_size=12, color=BLUE,
                align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text, subtitle=None):
    """슬라이드 상단 제목 영역"""
    # 제목 텍스트
    add_textbox(slide, title_text,
                ML, MT, CW, Cm(1.4),
                font_name=F_BOLD, font_size=28, color=NAVY)
    # 하단 구분선 (얇은 직사각형)
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        ML, MT + Cm(1.45), CW, Cm(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    if subtitle:
        add_textbox(slide, subtitle,
                    ML, MT + Cm(1.6), CW, Cm(0.8),
                    font_name=F_BODY, font_size=16, color=DGRAY)


def make_table(slide, headers, rows, x, y, col_widths,
               row_height=Cm(0.75), font_size=14):
    """표 생성. headers: 헤더 텍스트 리스트, rows: 데이터 행 리스트"""
    from pptx.util import Pt
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        x, y,
        sum(col_widths), row_height * n_rows
    ).table

    # 열 너비
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # 행 높이
    for r in range(n_rows):
        tbl.rows[r].height = row_height

    def set_cell(cell, text, bg, fg, bold=False, center=False, fsize=None):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name  = F_BOLD if bold else F_BODY
        run.font.size  = Pt(fsize or font_size)
        run.font.bold  = bold
        run.font.color.rgb = fg

    # 헤더 행
    for c, h in enumerate(headers):
        set_cell(tbl.cell(0, c), h, BLUE, WHITE, bold=True, center=True)

    # 데이터 행
    for r, row in enumerate(rows):
        row_bg = BG if r % 2 == 0 else LGRAY
        for c, val in enumerate(row):
            if isinstance(val, dict):
                set_cell(tbl.cell(r+1, c), val['text'], row_bg,
                         val.get('color', NAVY),
                         bold=val.get('bold', False),
                         center=val.get('center', False))
            else:
                set_cell(tbl.cell(r+1, c), str(val), row_bg, NAVY)

    return tbl


def add_info_box(slide, title, lines, x, y, w,
                 font_size=15, title_size=15):
    """파란 헤더 + 내용 박스"""
    line_h = Cm(0.62)
    title_h = Cm(0.7)
    body_h  = line_h * len(lines) + Cm(0.3)
    total_h = title_h + body_h

    # 헤더 박스
    hdr = slide.shapes.add_shape(1, x, y, w, title_h)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = BLUE
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.name = F_BOLD; run.font.size = Pt(title_size)
    run.font.bold = True; run.font.color.rgb = WHITE

    # 본문 박스
    body = slide.shapes.add_shape(1, x, y + title_h, w, body_h)
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
    body.line.fill.background()
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name  = F_BODY
        run.font.size  = Pt(font_size)
        run.font.color.rgb = NAVY

    return total_h


# ════════════════════════════════════════════════════════
# 슬라이드 생성
# ════════════════════════════════════════════════════════

prs   = new_prs()
blank = prs.slide_layouts[6]  # 빈 레이아웃
TOTAL = 18

# ── 슬라이드 1: 표지 ─────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)

# 상단 파란 띠
bar = sl.shapes.add_shape(1, 0, 0, W, Cm(0.5))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

# 하단 네이비 띠
bot = sl.shapes.add_shape(1, 0, H - Cm(2.2), W, Cm(2.2))
bot.fill.solid(); bot.fill.fore_color.rgb = NAVY
bot.line.fill.background()

add_textbox(sl, '캡스톤 디자인  |  6주차 연구 발표',
            0, Cm(2.5), W, Cm(1.0),
            font_name=F_MED, font_size=20, color=DGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, 'BTC 동적 그리드 트레이딩',
            0, Cm(4.0), W, Cm(2.0),
            font_name=F_TITLE, font_size=40, bold=False,
            color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(sl, 'RL 에이전트',
            0, Cm(5.8), W, Cm(2.0),
            font_name=F_BOLD, font_size=40,
            color=BLUE, align=PP_ALIGN.CENTER)

# 구분선
ln = sl.shapes.add_shape(1, Cm(10), Cm(8.2), Cm(13.87), Cm(0.06))
ln.fill.solid(); ln.fill.fore_color.rgb = LGRAY
ln.line.fill.background()

add_textbox(sl, '2026년 4월',
            0, Cm(8.6), W, Cm(0.8),
            font_name=F_BODY, font_size=16, color=DGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, '21101215  이찬희',
            0, H - Cm(1.6), W, Cm(1.0),
            font_name=F_MED, font_size=18, color=WHITE, align=PP_ALIGN.CENTER)


# ── 슬라이드 2: 목차 ─────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, 'Contents')
add_slide_number(sl, 2, TOTAL)

contents = [
    ('01', '연구 개요'),
    ('02', '데이터'),
    ('03', 'MDP 설계'),
    ('04', '매수/매도 실행 로직'),
    ('05', '베이스라인 전략 및 평가 결과'),
    ('06', 'PPO 학습 결과 및 분석'),
    ('07', '한계 및 향후 계획'),
]

col1_x = ML
col2_x = ML + Cm(1.8)
start_y = MT + Cm(2.2)
gap     = Cm(1.6)

for i, (num, title) in enumerate(contents):
    y = start_y + gap * i
    add_textbox(sl, num, col1_x, y, Cm(1.6), Cm(1.2),
                font_name=F_BOLD, font_size=22, color=BLUE)
    add_textbox(sl, title, col2_x, y, Cm(20), Cm(1.2),
                font_name=F_TITLE, font_size=22, color=NAVY)


def section_slide(prs, num, title, subtitle=None):
    """섹션 구분 슬라이드 (네이비 배경)"""
    sl = prs.slides.add_slide(blank)
    set_bg(sl, NAVY)

    # 좌측 강조 바
    bar = sl.shapes.add_shape(1, 0, 0, Cm(0.5), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_textbox(sl, num,
                Cm(1.5), Cm(6.5), Cm(5), Cm(2.0),
                font_name=F_BOLD, font_size=48, color=BLUE, align=PP_ALIGN.LEFT)
    add_textbox(sl, title,
                Cm(1.5), Cm(8.8), W - Cm(3), Cm(2.0),
                font_name=F_TITLE, font_size=36, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(sl, subtitle,
                    Cm(1.5), Cm(11.0), W - Cm(3), Cm(1.0),
                    font_name=F_BODY, font_size=18, color=LGRAY, align=PP_ALIGN.LEFT)
    return sl


# ── 슬라이드 3: 섹션1 구분 ───────────────────────────
sl = section_slide(prs, '01', '연구 개요')
add_slide_number(sl, 3, TOTAL)


# ── 슬라이드 4: 연구 주제 & 질문 ─────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '연구 주제 & 질문')
add_slide_number(sl, 4, TOTAL)

cy = MT + Cm(2.0)
add_info_box(sl,
    '연구 주제',
    ['가격 방향 예측 없이 시장 상태(변동성·가격 수준)와 포지션 상태(손익·보유량·여력)만을',
     'state로 사용하는 PPO 기반 강화학습 에이전트의 비트코인 동적 그리드 트레이딩 전략 학습 및 행동 패턴 분석'],
    ML, cy, CW, font_size=16)

cy += Cm(2.8)
add_info_box(sl,
    '주 질문',
    ['시장 상태와 포지션 상태에 반응하여 그리드 간격과 익절 목표를 동적으로 결정하는',
     'PPO 에이전트가 비트코인 시장에서 고정 그리드 전략 대비 Sharpe Ratio 기준 우위를 보이는가?'],
    ML, cy, CW, font_size=16)

cy += Cm(2.8)
add_info_box(sl,
    '부 질문',
    ['학습된 에이전트는 어떤 시장 상태(변동성 수준, 가격 수준)에서',
     '어떤 그리드 간격(aggressiveness)과 익절 목표(profit_target)를 선택하는가?'],
    ML, cy, CW, font_size=16)


# ── 슬라이드 5: 섹션2 구분 ───────────────────────────
sl = section_slide(prs, '02', '데이터')
add_slide_number(sl, 5, TOTAL)


# ── 슬라이드 6: 데이터 구성 ──────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '데이터 구성')
add_slide_number(sl, 6, TOTAL)

# 분할 표
make_table(sl,
    ['구분', '기간', '용도'],
    [
        ['Train',      '2020.01 ~ 2022.12  (25,916봉)', 'PPO 학습 — 상승·하락·횡보 레짐 혼재'],
        ['Validation', '2023.01 ~ 2023.12  (8,736봉)',  '학습 중 평가 (Val Sharpe 모니터링)'],
        ['Test',       '2024.01 ~ 현재',                {'text':'최종 평가 전까지 완전 봉인', 'color': RGBColor(0xC0,0x00,0x00), 'bold': True}],
    ],
    ML, MT + Cm(2.0),
    [Cm(3.2), Cm(8.0), Cm(12.0)],
    row_height=Cm(0.9)
)

# head 5행 표 (축약: timestamp, close, log_price, zscore_log_price, zscore_volatility)
add_textbox(sl, '처리된 데이터 구조 (Train 셋 상위 5행, 주요 컬럼)',
            ML, MT + Cm(5.6), CW, Cm(0.7),
            font_name=F_BOLD, font_size=16, color=NAVY)

make_table(sl,
    ['timestamp (UTC)', 'close', 'log_price', 'zscore_log_price', 'zscore_volatility'],
    [
        ['2020-01-14 22:00', '8,735.87', '0.0716', '1.015', '0.194'],
        ['2020-01-14 23:00', '8,810.01', '0.0796', '1.300', '0.059'],
        ['2020-01-15 00:00', '8,772.86', '0.0751', '1.184', '-0.007'],
        ['2020-01-15 01:00', '8,777.99', '0.0754', '1.229', '-0.056'],
        ['2020-01-15 02:00', '8,813.83', '0.0791', '1.386', '-0.211'],
    ],
    ML, MT + Cm(6.5),
    [Cm(5.5), Cm(3.2), Cm(3.2), Cm(4.4), Cm(4.4)],
    row_height=Cm(0.75)
)


# ── 슬라이드 7: 전처리 공식 ──────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '전처리 — 파생 지표 & 정규화')
add_slide_number(sl, 7, TOTAL)

cy = MT + Cm(2.0)
h1 = add_info_box(sl, 'log_price',
    ['log_price = log(close / close.rolling(168).mean())',
     '→ 현재 가격이 최근 1주일(168봉) 평균 대비 얼마나 위/아래인가'],
    ML, cy, CW, font_size=16)

cy += h1 + Cm(0.4)
h2 = add_info_box(sl, 'volatility_raw  (ATR 기반)',
    ['volatility_raw = ATR(168) / price',
     'ATR = mean(True Range) over 168봉',
     'True Range = max(high - low,  |high - prev_close|,  |low - prev_close|)'],
    ML, cy, CW, font_size=16)

cy += h2 + Cm(0.4)
add_info_box(sl, 'Rolling z-score 정규화  (window = 168)',
    ['zscore(x, t) = ( x_t  -  mean(x_{t-167:t}) )  /  std(x_{t-167:t})',
     '→ 두 변수 모두 정규화 후 zscore_log_price, zscore_volatility 컬럼으로 저장'],
    ML, cy, CW, font_size=16)


# ── 슬라이드 8: 섹션3 구분 ───────────────────────────
sl = section_slide(prs, '03', 'MDP 설계')
add_slide_number(sl, 8, TOTAL)


# ── 슬라이드 9: State & Action ────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, 'State Space  &  Action Space')
add_slide_number(sl, 9, TOTAL)

add_textbox(sl, 'State  (5차원 연속, rolling z-score 정규화)',
            ML, MT + Cm(1.9), CW, Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)

make_table(sl,
    ['차원', '변수명', '의미', '미보유 시'],
    [
        ['[0]', 'log_price',           '현재가 / 7일 이동평균 비율 (로그)',         '그대로'],
        ['[1]', 'divergence',          '현재가와 평균매수단가의 괴리율',              '아래 참조'],
        ['[2]', 'holdings_value_ratio','보유 BTC 평가액 / 초기 자본',               '0.0'],
        ['[3]', 'cash_ratio',          '현재 현금 / 초기 자본',                    '그대로'],
        ['[4]', 'volatility',          'ATR(168) / 현재가 (상대 변동성)',           '그대로'],
    ],
    ML, MT + Cm(2.7),
    [Cm(1.8), Cm(5.8), Cm(11.0), Cm(4.1)],
    row_height=Cm(0.78)
)

add_textbox(sl, 'Action  (2차원 연속 [0, 1]²)',
            ML, MT + Cm(7.4), CW, Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)

make_table(sl,
    ['차원', '변수명', '역할'],
    [
        ['[0]', 'aggressiveness', '매수 주문 간격 결정 — 클수록 현재가에서 멀리 주문  (0.01% ~ 5%)'],
        ['[1]', 'profit_target',  '매도 주문 간격 결정 — 클수록 더 높은 목표 수익률 설정  (0.01% ~ 15%)'],
    ],
    ML, MT + Cm(8.2),
    [Cm(1.8), Cm(5.8), Cm(15.1)],
    row_height=Cm(0.85)
)

add_textbox(sl, '※ divergence: 보유 중→avg_price 기준 / 미보유+직전 사이클→last_avg_price 기준 / 이력 없음→0.0',
            ML, H - Cm(1.5), CW, Cm(0.7),
            font_name=F_BODY, font_size=13, color=DGRAY)


# ── 슬라이드 10: Reward ───────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '보상 함수  (Reward)')
add_slide_number(sl, 10, TOTAL)

cy = MT + Cm(2.0)
add_info_box(sl, '보상 공식',
    ['step_reward  =  (equity_t  -  equity_{t-1})  /  start_capital  -  fee_rate × n_trades',
     '',
     '  equity_t    :  현재 스텝 포트폴리오 총가치 (현금 + 보유 BTC 평가액)',
     '  fee_rate    :  0.05%  (Binance maker fee)',
     '  n_trades    :  해당 스텝에서 체결된 주문 수',
     '',
     '  → 사이클 종료 시 별도 보너스 없음. completed_cycles 리스트에 통계만 기록.'],
    ML, cy, CW, font_size=16)

cy += Cm(5.8)
add_info_box(sl, 'Sharpe Ratio 연율화',
    ['Sharpe  =  ( mean(r_t) / std(r_t) )  ×  √8760',
     '',
     '  r_t    :  매 시간봉 포트폴리오 수익률',
     '  √8760  :  연율화 계수  (1년 = 8,760시간봉)'],
    ML, cy, CW, font_size=16)


# ── 슬라이드 11: 섹션4 구분 ──────────────────────────
sl = section_slide(prs, '04', '매수/매도 실행 로직')
add_slide_number(sl, 11, TOTAL)


# ── 슬라이드 12: 주문 가격 & 크기 ────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '주문 가격 계산  &  예산 분할')
add_slide_number(sl, 12, TOTAL)

cy = MT + Cm(2.0)
add_info_box(sl, 'Action → 주문 가격 변환',
    ['[ 매수 ]  기준: 현재가(price)',
     '  buy_hi  = price × (1 - (0.0001 + agg × 0.05))    →  [0.01%,  5%] 아래',
     '  buy_lo  = price × (1 - (0.001  + agg × 0.10))    →  [0.10%, 10%] 아래',
     '',
     '[ 매도 ]  sell_lo / sell_hi 는 기준 가격이 의도적으로 다름',
     '  sell_lo = price     × (1 + (0.0001 + pt × 0.05))  →  시장 모멘텀 활용 (단기 소익)',
     '  sell_hi = avg_price × (1 + (0.001  + pt × 0.15))  →  원가 기반 수익 보호'],
    ML, cy, CW, font_size=15)

cy += Cm(5.6)
add_info_box(sl, '주문 크기  (n_splits 예산 분할)',
    ['cycle_slot_size = cycle_start_cash / n_splits  (=4)',
     'per_order_size  = cycle_slot_size  / n_buy_orders  (=2)',
     'threshold_btc   = cycle_slot_size  / price',
     '',
     '예시: $10,000 → slot=$2,500 / 주문당 $1,250  |  슬롯 소진 후 매수 완전 차단'],
    ML, cy, CW, font_size=15)


# ── 슬라이드 13: 체결 판단 & SELL 우선 ───────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '체결 판단  &  SELL 우선 원칙')
add_slide_number(sl, 13, TOTAL)

cy = MT + Cm(2.0)
add_info_box(sl, '체결 조건  (다음 봉 high / low 기준)',
    ['① next_high  >=  sell_lo   →  sell_lo 체결',
     '② next_high  >=  sell_hi   →  sell_hi 체결  (avg_price 기준)',
     '③ next_low   <=  buy_hi    →  buy_hi  체결',
     '④ next_low   <=  buy_lo    →  buy_lo  체결',
     '',
     '  같은 봉에서 buy & sell 동시 조건 충족 시 →  ①②를 먼저, ③④를 나중에 처리'],
    ML, cy, CW, font_size=16)

cy += Cm(5.2)
add_info_box(sl, 'SELL 우선 원칙의 이유',
    ['1. 같은 봉에서 매도 후 수익을 먼저 확정하고 재매수하는 것이 유리',
     '2. 동시 처리 시 포트폴리오 회계(현금 잔고, 평단가)가 불명확해지는 것을 방지'],
    ML, cy, CW, font_size=16)

cy += Cm(2.8)
add_info_box(sl, '매도 수량 결정  (threshold_btc)',
    ['holdings  ≤  threshold_btc   →  전량 청산  (사이클 종료)',
     'holdings  >  threshold_btc   →  holdings / n_splits  균등 분할 매도'],
    ML, cy, CW, font_size=16)


# ── 슬라이드 14: 섹션5 구분 ──────────────────────────
sl = section_slide(prs, '05', '베이스라인 전략 및 평가 결과')
add_slide_number(sl, 14, TOTAL)


# ── 슬라이드 15: 베이스라인 결과 ─────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '베이스라인 Train / Val 평가 결과')
add_slide_number(sl, 15, TOTAL)

add_textbox(sl, '★ PPO 목표:  Train Sharpe > 0.818  (ATR k=0.5)  /  Val Sharpe > 2.610  (Fixed Grid 1%)',
            ML, MT + Cm(1.9), CW, Cm(0.7),
            font_name=F_BOLD, font_size=15, color=ORANGE)

GREEN_C = RGBColor(0x37, 0x56, 0x23)
RED_C   = RGBColor(0xC0, 0x00, 0x00)

make_table(sl,
    ['전략', 'Train Sharpe', 'Val Sharpe', 'Train 수익률', 'Val 수익률', 'Val MDD'],
    [
        ['Buy-and-Hold',      '0.666',                                 '2.377',                                '+89.9%', '+150.2%', '21.7%'],
        ['Fixed Grid 1%',     '0.620',                                 {'text':'2.610 ★','color':GREEN_C,'bold':True,'center':True}, '+73.4%', '+43.2%', '10.8%'],
        ['Fixed Grid 2%',     '0.205',                                 '2.032',                                '-2.0%',  '+17.0%',  '7.7%'],
        ['Fixed Grid 5%',     '0.802',                                 '1.375',                                '+50.0%', '+2.5%',   '1.7%'],
        ['ATR Grid k=0.5',    {'text':'0.818 ★','color':GREEN_C,'bold':True,'center':True}, '1.118',          '+122.1%','2+4.7%',  '16.1%'],
        ['ATR Grid k=1.0',    '0.686',                                 '1.434',                                '+86.4%', '+39.8%',  '15.9%'],
        ['ATR Grid k=2.0',    {'text':'-0.104','color':RED_C,'center':True}, '1.948',                         '-46.5%', '+29.0%',  '9.4%'],
    ],
    ML, MT + Cm(2.8),
    [Cm(4.4), Cm(3.4), Cm(3.4), Cm(3.4), Cm(3.4), Cm(4.7)],
    row_height=Cm(0.82)
)

add_textbox(sl,
    'Train(2020-2022)과 Val(2023) 순위 역전 주목 — Val은 BTC $16k→$42k 강세장으로 Buy-and-Hold 유리.\n'
    'PPO 학습 환경은 Train이므로 Train Sharpe 기준 비교가 더 공정.',
    ML, H - Cm(2.0), CW, Cm(1.4),
    font_name=F_BODY, font_size=14, color=DGRAY)


# ── 슬라이드 16: 섹션6 구분 ──────────────────────────
sl = section_slide(prs, '06', 'PPO 학습 결과 및 분석')
add_slide_number(sl, 16, TOTAL)


# ── 슬라이드 17: exp001 vs exp002 ────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, 'PPO  exp001 vs exp002')
add_slide_number(sl, 17, TOTAL)

# 왼쪽: 하이퍼파라미터 표
add_textbox(sl, '하이퍼파라미터',
            ML, MT + Cm(1.9), Cm(11), Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)
make_table(sl,
    ['파라미터', 'exp001', 'exp002'],
    [
        ['learning_rate', '3e-4',      '1e-4  ↓'],
        ['n_steps',       '2,048',     '4,096  ↑'],
        ['batch_size',    '64',        '128  ↑'],
        ['ent_coef',      '0.01',      '0.005  ↓'],
        ['total_steps',   '1,000,000', '1,000,000'],
    ],
    ML, MT + Cm(2.7),
    [Cm(4.8), Cm(2.8), Cm(3.0)],
    row_height=Cm(0.8)
)

# 오른쪽: 결과 표
add_textbox(sl, '학습 결과',
            ML + Cm(11.2), MT + Cm(1.9), Cm(11), Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)
make_table(sl,
    ['지표', 'exp001', 'exp002'],
    [
        ['best Sharpe',           '0.795 (100k)',  '0.745 (50k)'],
        ['final Sharpe',          {'text':'-0.280','color':RED_C,'center':True}, '-0.095'],
        ['평균 Sharpe',            '0.271',         '0.368  ↑'],
        ['양수 비율',               '95.0%',         '90.0%'],
        ['Train 베이스라인 최고 대비', {'text':'미달','color':RED_C,'center':True}, {'text':'미달','color':RED_C,'center':True}],
    ],
    ML + Cm(11.2), MT + Cm(2.7),
    [Cm(5.5), Cm(2.6), Cm(2.6)],
    row_height=Cm(0.8)
)

add_textbox(sl,
    '해석: exp002는 안정성 개선(평균 Sharpe 0.271→0.368, final -0.280→-0.095).\n'
    '그러나 두 실험 모두 Train 베이스라인(0.818) 미달. 단순 lr·ent_coef 조정만으로는 한계.',
    ML, H - Cm(2.0), CW, Cm(1.4),
    font_name=F_BODY, font_size=14, color=DGRAY)


# ── 슬라이드 18: 섹션7 한계 & 향후 계획 ─────────────
sl = prs.slides.add_slide(blank)
set_bg(sl, BG)
add_title_bar(sl, '한계 및 향후 계획')
add_slide_number(sl, 18, TOTAL)

add_textbox(sl, '발견된 설계 이슈 (exp003 반영 예정)',
            ML, MT + Cm(1.9), CW, Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)

make_table(sl,
    ['이슈', '개선 방향'],
    [
        ['threshold_btc 분모 불명확',        'avg(sell_lo,sell_hi) → price 로 단순화'],
        ['균등 분할 매도 기준 불일치',          'holdings/n_sell_orders → holdings/n_splits (매수 대칭)'],
        ['지정가 체결 방식 실전 불일치',         '조건 충족 시 지정가 → 그 시점 시장가로 체결'],
        ['sell_lo/sell_hi 명칭 혼동',         'sell_market / sell_cost 로 변경 (기준 가격 명시)'],
    ],
    ML, MT + Cm(2.7),
    [Cm(8.0), Cm(14.7)],
    row_height=Cm(0.82)
)

add_textbox(sl, '하이퍼파라미터 튜닝 방향 (exp003)',
            ML, MT + Cm(7.2), CW, Cm(0.7),
            font_name=F_BOLD, font_size=16, color=BLUE)

make_table(sl,
    ['방법', '내용', '기대 효과'],
    [
        ['VecNormalize',  'SB3 내장 reward/obs 자동 정규화',    '약한 보상 신호 문제 해결 (근본 원인)'],
        ['LR Scheduling', '3e-4 → 1e-5  cosine decay',        '후반부 수렴 불안정 해소'],
        ['더 긴 학습',     '1M → 3M 스텝',                     '충분한 샘플 확보'],
        ['n_steps 증가',  '4,096 → 8,192',                    '긴 에피소드 대응 (크레딧 할당 개선)'],
    ],
    ML, MT + Cm(8.0),
    [Cm(4.2), Cm(7.0), Cm(11.5)],
    row_height=Cm(0.78)
)


# ── 저장 ─────────────────────────────────────────────
out = 'reports/semester1/week6_presentation.pptx'
os.makedirs('reports/semester1', exist_ok=True)
prs.save(out)
print(f'저장 완료: {out}')
